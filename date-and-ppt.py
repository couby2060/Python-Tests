import datetime

now = datetime.datetime.now()
date = now.strftime("%Y-%m-%d")


from pptx import Presentation
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

filename = (date+"_test.pptx")
directory = ("read-write-folder/files-to-write/")
prs.save(directory+filename)